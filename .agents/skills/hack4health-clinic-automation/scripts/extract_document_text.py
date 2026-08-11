import os
import sys
import argparse
import docx
from pptx import Presentation

# Handle console encoding for Windows
try:
    sys.stdout.reconfigure(encoding='utf-8')
except AttributeError:
    pass

def extract_docx(file_path):
    doc = docx.Document(file_path)
    output = []
    
    # We will traverse the document elements in order of appearance
    # to keep tables and paragraphs properly interleaved.
    body_elements = doc.element.body
    paragraphs = []
    tables = []
    
    # Track paragraph and table indices
    p_idx = 0
    t_idx = 0
    
    # Iterate through child elements of the body
    for child in body_elements:
        if child.tag.endswith('p'):
            # Paragraph
            p_text = doc.paragraphs[p_idx].text.strip()
            if p_text:
                output.append(p_text)
            p_idx += 1
        elif child.tag.endswith('tbl'):
            # Table
            table = doc.tables[t_idx]
            output.append("\n[TABLE]")
            
            # Format table as Markdown Table
            markdown_table = []
            for r_idx, row in enumerate(table.rows):
                row_cells = []
                for cell in row.cells:
                    # Clean text and replace newlines with space/break
                    cell_text = cell.text.strip().replace("\n", " ").replace("|", "\\|")
                    row_cells.append(cell_text)
                
                # Check for duplicates (docx sometimes replicates cells when merged)
                # To clean it, we can keep the cell text but check if adjacent columns are identical
                cleaned_cells = []
                for c in row_cells:
                    if not cleaned_cells or cleaned_cells[-1] != c:
                        cleaned_cells.append(c)
                
                markdown_table.append("| " + " | ".join(cleaned_cells) + " |")
                
                if r_idx == 0 and len(cleaned_cells) > 0:
                    markdown_table.append("| " + " | ".join(["---"] * len(cleaned_cells)) + " |")
                    
            output.append("\n".join(markdown_table))
            output.append("[/TABLE]\n")
            t_idx += 1
            
    # Fallback if interleaving missed anything
    if not output:
        for p in doc.paragraphs:
            if p.text.strip():
                output.append(p.text.strip())
        for t in doc.tables:
            output.append("\n[TABLE]")
            for row in t.rows:
                cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                output.append("| " + " | ".join(cells) + " |")
            output.append("[/TABLE]\n")
            
    return "\n\n".join(output)

def extract_pptx(file_path):
    prs = Presentation(file_path)
    output = []
    
    for i, slide in enumerate(prs.slides):
        output.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                output.append(shape.text.strip())
    return "\n\n".join(output)

def main():
    parser = argparse.ArgumentParser(description="Extract clean text/markdown from docx/pptx files")
    parser.add_argument("file_path", help="Path to docx or pptx file")
    parser.add_argument("-o", "--output", help="Output file path (saves text to file)")
    
    args = parser.parse_args()
    
    if not os.path.exists(args.file_path):
        print(f"Error: File not found at '{args.file_path}'", file=sys.stderr)
        sys.exit(1)
        
    ext = os.path.splitext(args.file_path)[1].lower()
    text = ""
    
    try:
        if ext == ".docx":
            text = extract_docx(args.file_path)
        elif ext == ".pptx":
            text = extract_pptx(args.file_path)
        else:
            print(f"Error: Unsupported file format '{ext}'. Must be .docx or .pptx", file=sys.stderr)
            sys.exit(1)
            
        if args.output:
            with open(args.output, "w", encoding="utf-8") as f:
                f.write(text)
            print(f"Successfully extracted and saved to '{args.output}'")
        else:
            print(text)
            
    except Exception as e:
        print(f"Error extracting file: {e}", file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()
