"""Authoritative deterministic clinic rules and document extraction."""

from __future__ import annotations

import io
import re
from datetime import date
from pathlib import Path
from typing import Any

from docx import Document
from pptx import Presentation

SUPPORTED_UPLOADS = {".docx", ".pptx", ".txt"}
MAX_UPLOAD_BYTES = 8 * 1024 * 1024


def calculate_age(dob: str, as_of: date | None = None) -> int | None:
    """Calculate age using full birth date; two-digit years pivot at 2026."""
    as_of = as_of or date.today()
    try:
        day, month, raw_year = (int(part) for part in dob.split("/"))
        year = raw_year
        if year < 100:
            year += 1900 if year > 26 else 2000
        born = date(year, month, day)
        return (
            as_of.year - born.year - ((as_of.month, as_of.day) < (born.month, born.day))
        )
    except (TypeError, ValueError):
        return None


def match_eligibility(code: str, dob: str, gender: str = "") -> dict[str, Any]:
    """Resolve a supplied TPA code without LLM inference."""
    insurer_code = code.strip().upper()
    age = calculate_age(dob, date(2026, 12, 31))
    result: dict[str, Any] = {
        "insurer_code": insurer_code,
        "age": age,
        "gender": gender,
        "package_code": "UNKNOWN",
        "package_name": "Unknown Package / Custom Referral",
        "covered_tests": [],
        "notes": "",
        "requires_manual_review": False,
    }
    packages = {
        "BLPDE": (
            "BLPDE_UW",
            "Bluepeak Underwriting Requirements",
            ["Full Medical Examination (No Paramedics)", "Anti-HIV Test (BP#HIV)"],
        ),
        "MRDEB": (
            "MRDEB_UW",
            "Meridian Life Medical Referral Letter",
            ["Chest X-Ray", "Meridian Life# 10 - HIV Antibody Test", "Treadmill ECG"],
        ),
        "EVWME": (
            "EVWME_VOUCHER",
            "Everwell Health Check-up Voucher",
            ["Medical Examination", "Lipid Profile", "UFEME (Urine)", "Resting ECG"],
        ),
        "EVWPA": (
            "EVWPA_UW",
            "Everwell Adult Medical Examination",
            ["Adult Medical Examination"],
        ),
        "NSTNBU": (
            "NSTNBU_UW",
            "Northstar Life Underwriting Follow Up",
            ["Two repeat Urine Examination and Microscopy on different days"],
        ),
    }
    if insurer_code == "BLPHS" and age is not None:
        if age < 40:
            package = (
                "WELL1",
                "BluePeak Essential Screen",
                [
                    "Complete History Taking",
                    "Complete Physical Examination",
                    "BMI and Fat Composition",
                    "Blood pressure measurement",
                    "Full Blood Count",
                    "Cholesterol screening",
                ],
            )
        elif age <= 59:
            package = (
                "WELL2",
                "BluePeak Comprehensive Screen",
                [
                    "Complete History Taking",
                    "Complete Physical Examination",
                    "BMI and Fat Composition",
                    "Blood pressure measurement",
                    "Full Blood Count",
                    "Cholesterol screening",
                    "Liver function screening",
                    "Kidney function screening",
                    "Thyroid function screening",
                    "Resting ECG",
                    "Medical Report Consultation",
                ],
            )
        else:
            package = (
                "WELL3",
                "BluePeak Executive Screen",
                [
                    "Complete History Taking",
                    "Complete Physical Examination",
                    "BMI and Fat Composition",
                    "Blood pressure measurement",
                    "Full Blood Count",
                    "Cholesterol screening",
                    "Liver function screening",
                    "Kidney function screening",
                    "Thyroid function screening",
                    "Resting ECG",
                    "Liver Cancer Marker",
                    "Medical Report Consultation",
                ],
            )
        result.update(
            package_code=package[0], package_name=package[1], covered_tests=package[2]
        )
    elif insurer_code == "MOL0199VME" and age is not None:
        package_code = "PEE225" if age <= 24 else "PEE226" if age <= 49 else "PEE224"
        result.update(
            package_code=package_code,
            package_name="Ministry of Learning Civil Service Medical Scheme",
            covered_tests=["Pre-employment medical examination"],
        )
    elif insurer_code in packages:
        package = packages[insurer_code]
        result.update(
            package_code=package[0], package_name=package[1], covered_tests=package[2]
        )
    else:
        result["requires_manual_review"] = True
        result["notes"] = (
            "Unknown code or invalid date of birth; staff verification required."
        )
    return result


def compare_coverage(requested: list[str], covered: list[str]) -> dict[str, list[str]]:
    normalized = {re.sub(r"\W+", "", item).casefold(): item for item in covered}
    uncovered = [
        item
        for item in requested
        if re.sub(r"\W+", "", item).casefold() not in normalized
    ]
    return {
        "covered": [item for item in requested if item not in uncovered],
        "uncovered": uncovered,
    }


def build_prefill(patient: dict[str, Any], form_type: str) -> dict[str, Any]:
    identifier = str(patient.get("NRIC/FIN/Passport Number", ""))
    id_type = (
        "NRIC/FIN"
        if len(identifier) == 9 and identifier[:1].upper() in "STFGM"
        else "Passport"
    )
    allergy = str(patient.get("Drug Allergy", "")).strip()
    has_allergy = bool(allergy and allergy.casefold() not in {"none", "nil"})
    payload: dict[str, Any] = {
        "Name": patient.get("Full Name", ""),
        "Select One": id_type,
        "NRIC/FIN no.": identifier if id_type == "NRIC/FIN" else "",
        "Passport": identifier if id_type == "Passport" else "",
        "Date of Birth": patient.get("Date of Birth (DD/MM/YY)", ""),
        "Email Address": patient.get("Email", ""),
        "Country code": "+65",
        "Phone Number": str(
            patient.get("Contact - Mobile")
            or patient.get("Contact - Home")
            or patient.get("Contact - Office")
            or ""
        ),
        "Address": patient.get("Address", ""),
        "Postal Code": str(patient.get("Postal Code", "")).zfill(6),
        "Gender": "Male"
        if str(patient.get("Sex", "")).casefold() in {"m", "male"}
        else "Female",
        "Do you have any drug allergies?": "Yes" if has_allergy else "No",
    }
    if has_allergy:
        payload["Please provide name(s) of the drug(s)"] = allergy
    if form_type == "occupational":
        payload.update(
            {
                "Occupational Health Screening Type": ["Pre/Re Employment"],
                "Health Screening Location": "Parkway Shenton Medical Clinic (Republic Plaza)",
            }
        )
    else:
        payload.update(
            {
                "Health Screening Provider": "Parkway Shenton Medical Clinic",
                "Health Screening Location": "Republic Plaza",
            }
        )
    return payload


def extract_document(filename: str, content: bytes) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix not in SUPPORTED_UPLOADS:
        raise ValueError("Unsupported file type")
    if not content or len(content) > MAX_UPLOAD_BYTES:
        raise ValueError("File must be non-empty and no larger than 8 MB")
    if suffix == ".txt":
        return content.decode("utf-8")
    if suffix == ".docx":
        document = Document(io.BytesIO(content))
        blocks = [
            paragraph.text.strip()
            for paragraph in document.paragraphs
            if paragraph.text.strip()
        ]
        blocks.extend(
            " | ".join(cell.text.strip() for cell in row.cells)
            for table in document.tables
            for row in table.rows
        )
        return "\n".join(blocks)
    presentation = Presentation(io.BytesIO(content))
    return "\n".join(
        shape.text.strip()
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text.strip()
    )
